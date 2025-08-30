from flask import Flask, render_template, request, send_file, jsonify
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import PP_PLACEHOLDER
import re
import os
from dotenv import load_dotenv
from openai import OpenAI
import uuid
import io
import httpx

# Load environment variables
load_dotenv()

# Provider configuration
PROVIDER_CONFIG = {
    'gemini': {
        'base_url': "https://generativelanguage.googleapis.com/v1beta/",
        'models': {
            'gemini-2.5-flash': 'gemini-2.5-flash',
            'gemini-1.5-pro': 'gemini-1.5-pro',
            'gemini-1.5-flash': 'gemini-1.5-flash'
        }
    },
    'openai': {
        'base_url': "https://api.openai.com/v1/",
        'models': {
            'gpt-4': 'gpt-4',
            'gpt-4-turbo': 'gpt-4-turbo',
            'gpt-3.5-turbo': 'gpt-3.5-turbo'
        }
    },
    'claude': {
        'base_url': "https://api.anthropic.com/v1/",
        'models': {
            'claude-3-5-sonnet': 'claude-3-5-sonnet-20240620',
            'claude-3-opus': 'claude-3-opus-20240229',
            'claude-3-sonnet': 'claude-3-sonnet-20240229'
        }
    }
}

# Initialize client as None, will be set dynamically
client = None

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = '/tmp/uploads'  # Use /tmp for Vercel compatibility

if not os.path.exists(app.config['UPLOAD_FOLDER']):
    os.makedirs(app.config['UPLOAD_FOLDER'])

def remove_slide(prs, slide_index):
    """Remove a slide by its index."""
    slides = prs.slides
    slide_id_lst = prs.slides._sldIdLst
    slide_id_lst.remove(slide_id_lst[slide_index])

def get_layout(prs, name, available_layouts):
    """Dynamically map a layout name to an index, with intelligent fallbacks."""
    for idx, layout_name in available_layouts.items():
        if name.upper() == layout_name.upper():
            return prs.slide_layouts[idx]
    for idx, layout_name in available_layouts.items():
        if name.lower() in layout_name.lower():
            return prs.slide_layouts[idx]
    name_lower = name.lower()
    if "title" in name_lower and "section" in name_lower:
        for idx, layout_name in available_layouts.items():
            if "SECTION" in layout_name.upper() or "TITLE" in layout_name.upper():
                return prs.slide_layouts[idx]
    elif "title" in name_lower:
        for idx, layout_name in available_layouts.items():
            if "TITLE" in layout_name.upper() and "SECTION" not in layout_name.upper():
                return prs.slide_layouts[idx]
    elif "section" in name_lower:
        for idx, layout_name in available_layouts.items():
            if "SECTION" in layout_name.upper():
                return prs.slide_layouts[idx]
    elif "blank" in name_lower:
        for idx, layout_name in available_layouts.items():
            if "BLANK" in layout_name.upper():
                return prs.slide_layouts[idx]
    for idx, layout_name in available_layouts.items():
        if "TITLE_AND_BODY" in layout_name.upper():
            return prs.slide_layouts[idx]
    return prs.slide_layouts[0]

def estimate_placeholder_capacity(shape):
    """Estimate the capacity of a placeholder based on its size."""
    if not hasattr(shape, 'width') or not hasattr(shape, 'height'):
        return 10
    width_inches = shape.width / 914400
    height_inches = shape.height / 914400
    area = width_inches * height_inches
    if area < 5:
        return 3
    elif area < 10:
        return 5
    elif area < 20:
        return 8
    elif area < 30:
        return 12
    else:
        return 15

def get_layout_constraints(prs, available_layouts):
    """Get layout constraints for better content generation."""
    constraints = {}
    for layout_idx, layout_name in available_layouts.items():
        layout = prs.slide_layouts[layout_idx]
        layout_info = {
            "title_capacity": 10,
            "body_capacity": 10,
            "has_body": False
        }
        for shape in layout.placeholders:
            if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                layout_info["title_capacity"] = min(8, estimate_placeholder_capacity(shape))
            elif shape.placeholder_format.type == PP_PLACEHOLDER.BODY:
                layout_info["body_capacity"] = estimate_placeholder_capacity(shape)
                layout_info["has_body"] = True
            elif shape.placeholder_format.type == PP_PLACEHOLDER.OBJECT:
                if shape.placeholder_format.idx != 0:
                    layout_info["body_capacity"] = max(
                        layout_info["body_capacity"],
                        estimate_placeholder_capacity(shape)
                    )
                    layout_info["has_body"] = True
        constraints[layout_name] = layout_info
    return constraints

def generate_markdown(content, available_layout_names, constraints_text, provider_name, model_name, api_key, custom_base_url=None, proxies=None):
    """Generate Markdown using the selected provider and model via OpenAI-compatible interface."""
    # Set up the client based on the selected provider
    global client
    provider = PROVIDER_CONFIG.get(provider_name)
    if not provider:
        raise ValueError(f"Unsupported provider: {provider_name}")
    
    # Use custom base URL if provided, otherwise use provider default
    base_url = custom_base_url if custom_base_url else provider['base_url']
    
    # Create a custom httpx client with proxy configuration if proxies are provided
    http_client = httpx.Client(proxies=proxies) if proxies else None
    
    client = OpenAI(
        api_key=api_key,
        base_url=base_url,
        http_client=http_client
    )
    
    # Get the model identifier for the selected provider
    # For custom models, use the model_name directly
    if custom_base_url:
        model_id = model_name
    else:
        model_id = provider['models'].get(model_name, list(provider['models'].values())[0])
    
    system_prompt = f"""
You are an expert presentation designer creating structured Markdown for PowerPoint slides. Follow these rules:

1. Structure with `##` for slide titles and `-` for bullet points
2. Use layout comments like `<!-- Layout: TITLE -->` before each slide
3. For TITLE slides: Only include a title, no content
4. For SECTION_HEADER slides: Only include a title, no content
5. For content slides: Include a title and bullet points
6. Preserve bullet point hierarchy using indentation (2 spaces per level)
7. Use markdown formatting for emphasis:
   - Use **bold** for important terms, key concepts, or section headers within bullet points
   - Apply bold formatting sparingly for maximum impact
   - Do NOT use *italic* or other formatting
8. Keep content within space constraints:
   - Limit titles to the specified word count for each layout
   - Limit content slides to the specified bullet point count for each layout
   - If you have more content, split it into multiple slides with clear sub-topics
   - Keep bullet points concise (1 short sentence each)
9. Do NOT include markdown tables - describe them in plain text instead as bullet points or paragraphs
10. When you need to present tabular data, convert it to bullet points with clear labels
11. Do NOT include placeholder text like "List of Americas partner and non-partner teams"
12. Do NOT include text that says "(Table format would be used here)"
13. Do NOT include any text that indicates missing content or placeholders
14. When presenting data that would normally be in a table, format it as a series of bullet points with clear labels
15. Use ONLY these EXACT layout names from the available layouts:
{chr(10).join([f"    - {name}" for name in available_layout_names])}

Layout space constraints:
{constraints_text}

Example output:
<!-- Layout: {available_layout_names[0] if available_layout_names else 'TITLE'} -->
## Product Launch

<!-- Layout: {available_layout_names[1] if len(available_layout_names) > 1 else 'SECTION_HEADER'} -->
## Product Overview

<!-- Layout: {available_layout_names[2] if len(available_layout_names) > 2 else 'TITLE_AND_CONTENT'} -->
## Key Features
- **Feature 1**: Brief description
- **Feature 2**: Brief description
  - Sub-point for feature 2
- **Feature 3**: Brief description

Example of tabular data converted to bullet points:
## Team Information
- **Americas Partner Team**: Description of the team's responsibilities
- **Americas Non-Partner Team**: Description of the team's responsibilities
- **EMEA Team**: Description of the team's responsibilities
"""
    try:
        response = client.chat.completions.create(
            model=model_id,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": f"User-provided content:\n{content}"}
            ],
            temperature=0.5
        )
        return response.choices[0].message.content
    except Exception as e:
        # Provide more specific error information
        error_msg = f"API Error: {str(e)}"
        if custom_base_url:
            error_msg += f" (Using custom base URL: {custom_base_url})"
        raise Exception(error_msg)

def parse_markdown(md_content):
    """Parse Markdown content into structured slides."""
    slides = []
    current_slide = {"layout": None, "title": None, "content": []}
    lines = md_content.splitlines()
    for line in lines:
        line = line.rstrip()
        if not line.strip():
            continue
        if line.startswith("<!-- Layout:"):
            if current_slide["title"] or current_slide["content"]:
                slides.append(current_slide)
                current_slide = {"layout": None, "title": None, "content": []}
            layout_name = line.split("<!-- Layout:")[1].split("-->")[0].strip()
            current_slide["layout"] = layout_name
        elif line.startswith("## "):
            current_slide["title"] = line[3:].strip()
        elif line.startswith("- ") or re.match(r"^\s+- ", line):
            current_slide["content"].append(line)
    if current_slide["title"] or current_slide["content"]:
        slides.append(current_slide)
    return slides

def add_formatted_content(shape, content_lines):
    """Add properly formatted content to a shape with markdown formatting support."""
    if not content_lines:
        return
    text_frame = shape.text_frame
    text_frame.clear()
    for i, line in enumerate(content_lines):
        indent_level = 0
        stripped_line = line
        # Handle different types of content
        while stripped_line.startswith("  "):
            indent_level += 1
            stripped_line = stripped_line[2:]
        if stripped_line.startswith("- "):
            stripped_line = stripped_line[2:]
        # Add content as paragraphs
        if i == 0:
            p = text_frame.paragraphs[0]
            p.level = indent_level
        else:
            p = text_frame.add_paragraph()
            p.level = indent_level
        process_markdown_formatting(p, stripped_line)
        if hasattr(p, 'font') and p.font:
            p.font.size = Pt(18)

def process_markdown_formatting(paragraph, text):
    """Process markdown formatting and apply it to the paragraph."""
    if not text:
        return
        
    # Handle bold formatting (**text** or __text__)
    # Split text by both ** and __ patterns
    parts = re.split(r'(\*\*.*?\*\*|__.*?__)', text)
    
    for part in parts:
        if not part:
            continue
        # Check for **bold** formatting
        if part.startswith('**') and part.endswith('**') and len(part) > 4:
            run = paragraph.add_run()
            run.text = part[2:-2]  # Remove the ** markers
            run.font.bold = True
        # Check for __bold__ formatting
        elif part.startswith('__') and part.endswith('__') and len(part) > 4:
            run = paragraph.add_run()
            run.text = part[2:-2]  # Remove the __ markers
            run.font.bold = True
        else:
            # Regular text
            run = paragraph.add_run()
            run.text = part

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/generate', methods=['POST'])
def generate():
    if 'ppt_file' not in request.files:
        return "No file part", 400
    file = request.files['ppt_file']
    if file.filename == '':
        return "No selected file", 400
    if file:
        user_content = request.form['user_content']
        # Get provider, model, and API key from form
        provider_name = request.form.get('provider', 'gemini')
        model_name = request.form.get('model', 'gemini-2.5-flash')
        api_key = request.form.get('api_key')
        
        # Get custom model parameters if provided
        custom_base_url = request.form.get('custom_base_url')
        custom_model_name = request.form.get('custom_model_name')
        
        # Use custom model if provided
        if custom_base_url and custom_model_name:
            # Validate custom base URL format
            if not custom_base_url.startswith(('http://', 'https://')):
                return "Custom base URL must start with 'http://' or 'https://'", 400
            # Ensure URL ends with a slash
            if not custom_base_url.endswith('/'):
                custom_base_url += '/'
                
            provider_name = 'openai'  # Custom models are OpenAI-compatible
            model_name = custom_model_name
        
        # Validate API key
        if not api_key:
            return "API key is required", 400
        
        upload_folder = os.path.join(app.config['UPLOAD_FOLDER'], str(uuid.uuid4()))
        os.makedirs(upload_folder, exist_ok=True)
        
        template_path = os.path.join(upload_folder, file.filename)
        file.save(template_path)
        
        prs = Presentation(template_path)
        original_slide_count = len(prs.slides)
        
        available_layouts = {i: layout.name for i, layout in enumerate(prs.slide_layouts)}
        available_layout_names = list(available_layouts.values())
        
        layout_constraints = get_layout_constraints(prs, available_layouts)
        constraints_text = "\n".join([
            f"    - {name}: Title max {info['title_capacity']} words, Content max {info['body_capacity']} bullet points"
            for name, info in layout_constraints.items()
        ])
        
        try:
            # Get proxy settings from request if available
            proxies = {
                "http://": request.form.get("http_proxy"),
                "https://": request.form.get("https_proxy"),
            } if request.form.get("http_proxy") or request.form.get("https_proxy") else None
            
            md_content = generate_markdown(user_content, available_layout_names, constraints_text, provider_name, model_name, api_key, custom_base_url, proxies)
        except Exception as e:
            return f"Error generating content: {str(e)}", 500
            
        slides = parse_markdown(md_content)
        
        for slide_data in slides:
            title = slide_data["title"]
            content_lines = slide_data["content"]
            layout_name = slide_data["layout"] or "TITLE_AND_CONTENT"
            layout = get_layout(prs, layout_name, available_layouts)
            slide = prs.slides.add_slide(layout)
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = title or "Untitled"
            
            layout_upper = layout_name.upper()
            is_content_slide = not (
                layout_upper == "TITLE" or
                layout_upper == "SECTION_HEADER" or
                "SECTION HEADER" in layout_upper or
                ("TITLE" in layout_upper and "SECTION" in layout_upper and not any(c in layout_upper for c in ["CONTENT", "BODY"])) or
                "BLANK" in layout_upper
            )
            
            if content_lines and is_content_slide:
                body_shape = None
                for shape in slide.placeholders:
                    if shape.placeholder_format.type == PP_PLACEHOLDER.BODY:
                        body_shape = shape
                        break
                if not body_shape:
                    for shape in slide.placeholders:
                        shape_name = getattr(shape, 'name', '').upper()
                        if any(keyword in shape_name for keyword in ['CONTENT', 'BODY', 'TEXT']):
                            body_shape = shape
                            break
                if not body_shape:
                    for shape in slide.placeholders:
                        if hasattr(shape, 'placeholder_format') and shape.placeholder_format.idx != 0:
                            if hasattr(shape, 'text_frame'):
                                body_shape = shape
                                break
                if not body_shape:
                    for shape in slide.placeholders:
                        if shape != title_shape and hasattr(shape, 'text_frame'):
                            body_shape = shape
                            break
                if body_shape:
                    add_formatted_content(body_shape, content_lines)

        for i in reversed(range(original_slide_count)):
            remove_slide(prs, i)
        
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        
        return send_file(output, as_attachment=True, download_name='generated_presentation.pptx', mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')

# Vercel requires the app to be accessible as a variable named "app"
application = app

if __name__ == '__main__':
    app.run(debug=True)
